#!/usr/bin/env python3
"""根据技术分享大纲生成 PPTX。

输出：
- docs/tech-sharing-outline.pptx

设计目标：
- 明亮、柔和、带一点人文感
- 正文尽量可直接演讲
- 需要补素材的部分单独放到 appendix slides
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
OUTPUT_PATH = DOCS_DIR / "tech-sharing-outline.pptx"
GENERATED_DIR = DOCS_DIR / ".generated" / "tech-sharing-outline"

SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_FONT = "PingFang SC"
BODY_FONT = "PingFang SC"
SERIF_FONT = "Songti SC"

PALETTE = {
    "ink": "34424E",
    "muted": "6D7B84",
    "cream": "F8F4EC",
    "paper": "FFFDFC",
    "blush": "E8C6B8",
    "rose": "DDAE9D",
    "sage": "A8BEA3",
    "moss": "7D9B83",
    "clay": "B98068",
    "sand": "E8D9C1",
    "gold": "D9C38E",
    "teal": "8FAFB2",
    "plum": "7D6677",
    "dark": "3E4D58",
    "appendix": "F4ECE3",
    "line": "D7CCBF",
    "white": "FFFFFF",
}


@dataclass(frozen=True)
class Box:
    x: float
    y: float
    w: float
    h: float


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.replace("#", "")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def rgba_tuple(hex_color: str, alpha: int = 255) -> tuple[int, int, int, int]:
    hex_color = hex_color.replace("#", "")
    return (
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
        alpha,
    )


def make_blob_background(path: Path, base_hex: str, blobs: Sequence[tuple[int, int, int, str, int]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGBA", (1920, 1080), rgba_tuple(base_hex))
    overlay = Image.new("RGBA", (1920, 1080), (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    for x, y, radius, color_hex, alpha in blobs:
        draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=rgba_tuple(color_hex, alpha))
    overlay = overlay.filter(ImageFilter.GaussianBlur(90))
    img = Image.alpha_composite(img, overlay)

    grain = Image.new("RGBA", (1920, 1080), (255, 255, 255, 0))
    gdraw = ImageDraw.Draw(grain)
    for offset in range(0, 1920, 80):
        gdraw.line((offset, 0, offset - 220, 1080), fill=(255, 255, 255, 10), width=2)
    grain = grain.filter(ImageFilter.GaussianBlur(8))
    img = Image.alpha_composite(img, grain)
    img.convert("RGB").save(path)


def ensure_backgrounds() -> dict[str, Path]:
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    paths = {
        "cover": GENERATED_DIR / "cover-bg.png",
        "content": GENERATED_DIR / "content-bg.png",
        "warm": GENERATED_DIR / "warm-bg.png",
        "appendix": GENERATED_DIR / "appendix-bg.png",
        "dark": GENERATED_DIR / "dark-bg.png",
    }
    make_blob_background(
        paths["cover"],
        PALETTE["cream"],
        [
            (1650, 180, 300, PALETTE["blush"], 140),
            (1450, 720, 260, PALETTE["sage"], 110),
            (280, 180, 230, PALETTE["sand"], 120),
            (360, 900, 300, PALETTE["rose"], 70),
        ],
    )
    make_blob_background(
        paths["content"],
        PALETTE["cream"],
        [
            (1700, 120, 240, PALETTE["sand"], 90),
            (1550, 880, 260, PALETTE["sage"], 80),
            (240, 920, 300, PALETTE["blush"], 70),
            (300, 180, 220, PALETTE["teal"], 35),
        ],
    )
    make_blob_background(
        paths["warm"],
        PALETTE["appendix"],
        [
            (1660, 160, 270, PALETTE["rose"], 110),
            (1400, 860, 250, PALETTE["gold"], 85),
            (260, 180, 220, PALETTE["sage"], 65),
            (320, 860, 260, PALETTE["sand"], 70),
        ],
    )
    make_blob_background(
        paths["appendix"],
        PALETTE["appendix"],
        [
            (1570, 180, 280, PALETTE["rose"], 95),
            (1520, 880, 290, PALETTE["sage"], 80),
            (220, 860, 260, PALETTE["sand"], 95),
        ],
    )
    make_blob_background(
        paths["dark"],
        PALETTE["dark"],
        [
            (1750, 180, 260, PALETTE["clay"], 90),
            (280, 880, 320, PALETTE["teal"], 50),
            (1320, 860, 260, PALETTE["gold"], 45),
        ],
    )
    return paths


def add_background(slide, image_path: Path) -> None:
    slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(SLIDE_W), height=Inches(SLIDE_H))


def add_round_rect(slide, box: Box, fill_hex: str, line_hex: str | None = None, line_width: float = 1.25, transparency: int = 0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(box.x),
        Inches(box.y),
        Inches(box.w),
        Inches(box.h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.fill.transparency = transparency
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, box: Box, fill_hex: str, line_hex: str | None = None, line_width: float = 1.0, transparency: int = 0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(box.x),
        Inches(box.y),
        Inches(box.w),
        Inches(box.h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.fill.transparency = transparency
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text: str,
    box: Box,
    font_size: float,
    color_hex: str,
    *,
    font_name: str = BODY_FONT,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.06,
) -> None:
    shape = slide.shapes.add_textbox(Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)


def add_paragraphs(
    slide,
    items: Iterable[str],
    box: Box,
    *,
    font_size: float = 18,
    color_hex: str = PALETTE["ink"],
    bullet: bool = True,
    font_name: str = BODY_FONT,
    line_space: float = 1.15,
    margin: float = 0.08,
) -> None:
    shape = slide.shapes.add_textbox(Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_space
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"• {item}" if bullet else item
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = rgb(color_hex)


def add_badge(slide, text: str, box: Box, fill_hex: str, color_hex: str, *, font_size: float = 12) -> None:
    add_round_rect(slide, box, fill_hex, transparency=0)
    add_text(
        slide,
        text,
        box,
        font_size,
        color_hex,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.02,
    )


def add_title_block(slide, eyebrow: str, title: str, subtitle: str | None = None, *, dark: bool = False) -> None:
    text_color = PALETTE["paper"] if dark else PALETTE["muted"]
    title_color = PALETTE["paper"] if dark else PALETTE["ink"]
    add_badge(
        slide,
        eyebrow,
        Box(0.78, 0.5, 1.85, 0.34),
        PALETTE["sand"] if not dark else PALETTE["clay"],
        PALETTE["ink"] if not dark else PALETTE["paper"],
    )
    add_text(slide, title, Box(0.72, 0.92, 8.8, 0.9), 27, title_color, font_name=TITLE_FONT, bold=True)
    if subtitle:
        add_text(slide, subtitle, Box(0.76, 1.64, 9.4, 0.55), 14.5, text_color)


def add_footer(slide, page_no: int, label: str, *, dark: bool = False) -> None:
    color = PALETTE["paper"] if dark else PALETTE["muted"]
    add_text(slide, label, Box(0.78, 7.0, 3.8, 0.25), 10.5, color)
    add_text(slide, f"{page_no:02d}", Box(12.2, 6.94, 0.46, 0.28), 11, color, bold=True, align=PP_ALIGN.RIGHT)


def add_card_with_title(slide, title: str, box: Box, *, fill_hex: str, line_hex: str = PALETTE["line"], title_size: float = 18, body_items: Sequence[str] | None = None, body_size: float = 15.5) -> None:
    add_round_rect(slide, box, fill_hex, line_hex)
    add_text(slide, title, Box(box.x + 0.16, box.y + 0.14, box.w - 0.32, 0.35), title_size, PALETTE["ink"], bold=True)
    if body_items:
        add_paragraphs(
            slide,
            body_items,
            Box(box.x + 0.12, box.y + 0.5, box.w - 0.24, box.h - 0.58),
            font_size=body_size,
            color_hex=PALETTE["ink"],
        )


def add_quote_card(slide, quote: str, box: Box, *, fill_hex: str, color_hex: str = PALETTE["ink"], subtext: str | None = None) -> None:
    add_round_rect(slide, box, fill_hex, transparency=0)
    add_text(slide, "“", Box(box.x + 0.2, box.y + 0.18, 0.45, 0.4), 34, color_hex, font_name=SERIF_FONT, bold=True)
    quote_height = box.h - (0.98 if subtext else 0.62)
    add_text(slide, quote, Box(box.x + 0.55, box.y + 0.36, box.w - 0.92, quote_height), 22, color_hex, font_name=SERIF_FONT, bold=True)
    if subtext:
        add_text(slide, subtext, Box(box.x + 0.56, box.y + box.h - 0.42, box.w - 0.9, 0.24), 10.5, PALETTE["muted"])


def add_stat_card(slide, value: str, label: str, note: str, box: Box, *, fill_hex: str) -> None:
    add_round_rect(slide, box, fill_hex, line_hex=PALETTE["line"])
    add_text(slide, value, Box(box.x + 0.14, box.y + 0.12, box.w - 0.28, 0.55), 28, PALETTE["ink"], bold=True)
    add_text(slide, label, Box(box.x + 0.15, box.y + 0.68, box.w - 0.3, 0.28), 13.5, PALETTE["muted"], bold=True)
    add_text(slide, note, Box(box.x + 0.15, box.y + 1.0, box.w - 0.3, box.h - 1.05), 11.8, PALETTE["ink"])


def add_image_contain(slide, image_path: Path, box: Box) -> None:
    with Image.open(image_path) as img:
        aspect = img.width / img.height
    box_aspect = box.w / box.h
    if aspect > box_aspect:
        width = box.w
        height = width / aspect
        x = box.x
        y = box.y + (box.h - height) / 2
    else:
        height = box.h
        width = height * aspect
        x = box.x + (box.w - width) / 2
        y = box.y
    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(width), height=Inches(height))


def add_flow_box(slide, title: str, body: str, box: Box, fill_hex: str) -> None:
    add_round_rect(slide, box, fill_hex, line_hex=PALETTE["line"])
    add_text(slide, title, Box(box.x + 0.14, box.y + 0.12, box.w - 0.28, 0.3), 15.5, PALETTE["ink"], bold=True)
    add_text(slide, body, Box(box.x + 0.14, box.y + 0.46, box.w - 0.28, box.h - 0.58), 12.2, PALETTE["ink"])


def add_arrow(slide, x: float, y: float, w: float, h: float, fill_hex: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.fill.transparency = 15
    shape.line.fill.background()


def build_cover(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["cover"])
    add_badge(slide, "技术分享 | 60 分钟", Box(0.88, 0.56, 2.4, 0.36), PALETTE["sand"], PALETTE["ink"])
    add_text(
        slide,
        "怎么让 AI 帮你改代码，\n但不把项目改坏",
        Box(0.86, 1.25, 6.8, 1.5),
        29,
        PALETTE["ink"],
        font_name=TITLE_FONT,
        bold=True,
    )
    add_text(
        slide,
        "从 OpenClaw 的困境到 Nanobot/Ava 的 Harness 实践",
        Box(0.9, 2.95, 6.8, 0.42),
        16,
        PALETTE["muted"],
    )
    add_quote_card(
        slide,
        "我不是只会用 AI 写代码，我会给 AI 设计边界。",
        Box(0.9, 3.7, 5.95, 1.18),
        fill_hex=PALETTE["paper"],
    )
    add_text(slide, "主讲人：方壶", Box(0.92, 5.28, 2.6, 0.3), 13, PALETTE["ink"], bold=True)
    add_text(slide, "Nanobot/Ava · Harness Engineering · Spec Governance", Box(0.92, 5.58, 5.2, 0.3), 11.5, PALETTE["muted"])
    logo_box = Box(9.55, 1.2, 3.0, 2.1)
    add_round_rect(slide, Box(9.15, 0.9, 3.7, 2.8), PALETTE["paper"], transparency=12)
    add_image_contain(slide, ROOT / "nanobot_logo.png", logo_box)
    add_text(
        slide,
        "轻量 agent 框架 + sidecar harness",
        Box(9.35, 3.9, 3.4, 0.32),
        13,
        PALETTE["ink"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "用更小、更可控的系统，验证“AI 改 AI”的工程边界。",
        Box(9.25, 4.28, 3.6, 0.46),
        12.5,
        PALETTE["muted"],
        align=PP_ALIGN.CENTER,
    )


def build_core_message(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["warm"])
    add_title_block(slide, "核心信息", "全场希望只记住这一句话")
    add_quote_card(
        slide,
        "我不是只会用 AI 写代码，我会给 AI 设计边界。",
        Box(1.1, 1.8, 11.1, 2.2),
        fill_hex=PALETTE["paper"],
        subtext="把一个个人 AI 助手做成有边界、有状态、有验证、有生命周期的 agent harness 系统",
    )
    badges = [
        ("有边界", PALETTE["sand"]),
        ("有状态", PALETTE["sage"]),
        ("有验证", PALETTE["blush"]),
        ("有生命周期", PALETTE["gold"]),
    ]
    x = 1.3
    for text, fill in badges:
        add_badge(slide, text, Box(x, 4.55, 2.25, 0.44), fill, PALETTE["ink"], font_size=14)
        x += 2.6
    add_text(
        slide,
        "重点不是“AI 会不会写代码”，而是“我能不能把它放进一个可靠的工程回路里”。",
        Box(1.3, 5.42, 10.7, 0.44),
        15.2,
        PALETTE["ink"],
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 2, "Takeaway")


def build_openclaw(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["content"])
    add_title_block(slide, "一、开场", "OpenClaw 的爆火与“龙虾悖论”")
    add_card_with_title(
        slide,
        "从现象出发",
        Box(0.82, 1.82, 5.35, 2.5),
        fill_hex=PALETTE["paper"],
        body_items=[
            "现象级开源 Agent 产品，能接管键鼠、整理文件、填表、起草邮件。",
            "它让大众第一次直观看到“Agent 真能动手干活”。",
            "但能力越强，权限越大，攻击面也越大。",
        ],
    )
    add_card_with_title(
        slide,
        "龙虾悖论",
        Box(0.82, 4.52, 5.35, 1.75),
        fill_hex=PALETTE["sand"],
        body_items=[
            "你想让它做的事越多，给它的权限就越大；权限越大，安全风险就越高。",
            "这不是 OpenClaw 独有问题，Cursor、Claude Code、Codex 都会遇到。",
        ],
        body_size=14.5,
    )
    placeholder = Box(6.55, 1.8, 5.95, 4.52)
    add_round_rect(slide, placeholder, PALETTE["paper"], line_hex=PALETTE["rose"])
    add_badge(slide, "待补截图", Box(6.88, 2.02, 1.2, 0.32), PALETTE["rose"], PALETTE["paper"], font_size=11)
    add_text(slide, "OpenClaw 新闻 / GitHub 热度视觉 hook", Box(6.9, 2.48, 4.8, 0.32), 17, PALETTE["ink"], bold=True)
    add_text(
        slide,
        "建议你后续替换成一张“热度 + 风险”同时成立的截图：\n1. GitHub repo 首页或趋势页\n2. 或媒体报道封面 + 星标数\n3. 画面最好能一眼看出“很火”",
        Box(6.9, 2.96, 4.9, 1.1),
        13.2,
        PALETTE["ink"],
    )
    add_text(
        slide,
        "这里故意留空，不拿模糊截图硬顶。",
        Box(6.92, 5.66, 4.8, 0.28),
        12,
        PALETTE["muted"],
    )
    add_footer(slide, 3, "OpenClaw Hook")


def build_sources(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["content"])
    add_title_block(slide, "二、范式", "Harness Engineering：三个来源，一个共识", "不是讲概念，而是说明这套工程方法已经有共识、有案例、有抽象。")
    cards = [
        (
            "OpenAI 的启发",
            [
                "Humans steer. Agents execute.",
                "工程师不只写代码，还要设计环境、表达意图、搭反馈回路。",
                "规则最好落成仓库里的版本化文件。",
            ],
            PALETTE["sand"],
        ),
        (
            "Anthropic 的启发",
            [
                "长时程 Agent 需要 initializer、feature list、progress file、init.sh、E2E 验证。",
                "没有 harness，Agent 往往不能持续可靠地工作。",
            ],
            PALETTE["paper"],
        ),
        (
            "Learn Harness Engineering",
            [
                "把做法收束成五子系统：Instructions / State / Verification / Scope / Lifecycle。",
                "Harness 不让模型更聪明，它让输出更可靠。",
            ],
            PALETTE["appendix"],
        ),
    ]
    x = 0.86
    for title, body, fill in cards:
        add_card_with_title(slide, title, Box(x, 2.1, 3.96, 3.6), fill_hex=fill, body_items=body, body_size=15)
        x += 4.18
    add_text(
        slide,
        "共识：真正的差距，不是模型 IQ，而是有没有一套把它管住、续上、验对的 harness。",
        Box(1.05, 6.15, 11.1, 0.36),
        15.2,
        PALETTE["ink"],
        align=PP_ALIGN.CENTER,
        bold=True,
    )
    add_footer(slide, 4, "Three Sources")


def build_harness_five(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["warm"])
    add_title_block(slide, "二、范式", "Harness 的五个子系统")
    add_quote_card(
        slide,
        "模型决定写什么代码。Harness 治理何时、何地、怎么写。",
        Box(4.0, 2.6, 5.4, 1.55),
        fill_hex=PALETTE["paper"],
        subtext="这页对应的是方法抽象，不是空口定义。",
    )
    cards = [
        ("Instructions", "告诉 AI 做什么、不要做什么。", Box(0.92, 1.8, 2.7, 1.3), PALETTE["sand"]),
        ("State", "记住做到哪了、约束是什么、上下文如何回流。", Box(9.78, 1.8, 2.45, 1.3), PALETTE["sage"]),
        ("Verification", "不是它说完成了，而是系统能验证对不对。", Box(0.92, 4.45, 2.7, 1.3), PALETTE["blush"]),
        ("Scope", "让 AI 一次只在允许的边界里工作。", Box(9.78, 4.45, 2.45, 1.3), PALETTE["gold"]),
        ("Lifecycle", "初始化、续跑、清理、重启、重建。", Box(4.4, 5.2, 4.55, 1.12), PALETTE["paper"]),
    ]
    for title, body, box, fill in cards:
        add_card_with_title(slide, title, box, fill_hex=fill, body_items=[body], body_size=13.2)
    add_footer(slide, 5, "Five Systems")


def build_ava_mapping(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["content"])
    add_title_block(slide, "二、落地", "Ava 把五子系统真正落进了仓库")
    rows = [
        ("Instructions", "AGENTS.md / CLAUDE.md / ava/templates/TOOLS.md / .specanchor/global"),
        ("State", ".specanchor 三级 Spec / ava/agent/bg_tasks.py / session 持久化"),
        ("Verification", "tests/guardrails / .githooks/pre-commit / ava/console/ui_build.py"),
        ("Scope", "nanobot/ 禁改规则 / patch_map / strict schema gate / Task Spec"),
        ("Lifecycle", "python -m ava 启动 / patch apply / auto_continue / auto rebuild / supervisor-first restart"),
    ]
    y = 1.72
    for idx, (name, desc) in enumerate(rows):
        fill = [PALETTE["sand"], PALETTE["paper"], PALETTE["appendix"], PALETTE["paper"], PALETTE["appendix"]][idx]
        add_round_rect(slide, Box(0.86, y, 8.7, 0.78), fill, line_hex=PALETTE["line"])
        add_text(slide, name, Box(1.05, y + 0.14, 1.55, 0.22), 16.5, PALETTE["ink"], bold=True)
        add_text(slide, desc, Box(2.45, y + 0.12, 6.8, 0.4), 13.3, PALETTE["ink"])
        y += 0.92
    add_quote_card(
        slide,
        "这页最关键：不是“我知道 harness 是什么”，而是“我已经把它分层放到了代码里”。",
        Box(9.9, 1.95, 2.45, 2.95),
        fill_hex=PALETTE["paper"],
        subtext="State + Scope 里最特别的一层，就是 SpecAnchor。",
    )
    add_badge(slide, "重点页", Box(10.08, 5.2, 0.96, 0.3), PALETTE["clay"], PALETTE["paper"], font_size=11)
    add_text(slide, "建议这页慢一点讲，让听众先理解“概念”已经变成了“工程构件”。", Box(9.95, 5.56, 2.2, 0.78), 12.3, PALETTE["muted"])
    add_footer(slide, 6, "Ava Mapping")


def build_specanchor_role(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["warm"])
    add_title_block(slide, "二、边界", "SpecAnchor 在 Harness 里的角色")
    add_quote_card(
        slide,
        "SpecAnchor = 面向 Agent 开发的 spec governance layer。",
        Box(1.0, 1.75, 11.1, 1.15),
        fill_hex=PALETTE["paper"],
        subtext="它解决的是 spec 如何组织、加载、同步、回流，而不是替代所有 workflow 框架。",
    )
    add_card_with_title(
        slide,
        "它负责什么",
        Box(1.0, 3.25, 5.4, 2.55),
        fill_hex=PALETTE["paper"],
        body_items=[
            "Global / Module / Task 分层",
            "上下文按需加载",
            "Task 结论回流模块知识",
            "Spec 新鲜度与覆盖治理",
        ],
    )
    add_card_with_title(
        slide,
        "它不负责什么",
        Box(6.9, 3.25, 5.4, 2.55),
        fill_hex=PALETTE["appendix"],
        body_items=[
            "不替你定义所有需求的写法",
            "不独占 Plan Approved 这类 gate",
            "不替代 OpenSpec / Spec Kit / SDD-RIPER-ONE",
            "它更像 control plane，而不是万能工作流产品",
        ],
    )
    add_footer(slide, 7, "Spec Governance")


def build_nano_overview(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["content"])
    add_title_block(slide, "三、系统", "Nanobot/Ava：为什么这套系统适合做受控实验")
    add_card_with_title(
        slide,
        "为什么是 Nanobot",
        Box(0.86, 1.82, 4.15, 2.7),
        fill_hex=PALETTE["paper"],
        body_items=[
            "不是因为它“更酷”，而是更轻、更可读、更容易让个人开发者真正掌控。",
            "README 直接写着：Inspired by OpenClaw, 99% fewer lines。",
            "它足够小，小到可以被完整理解；又足够完整，能承载真实 agent 能力。",
        ],
        body_size=14.2,
    )
    add_card_with_title(
        slide,
        "Ava 的结构选择",
        Box(0.86, 4.75, 4.15, 1.55),
        fill_hex=PALETTE["sand"],
        body_items=[
            "不是 fork，是 sidecar。",
            "启动入口走 `python -m ava`，所有定制尽量收束在 `ava/`。",
        ],
        body_size=14,
    )
    image_card = Box(5.38, 1.72, 7.1, 4.68)
    add_round_rect(slide, image_card, PALETTE["paper"], line_hex=PALETTE["line"])
    add_image_contain(slide, ROOT / "nanobot_arch.png", Box(5.62, 2.0, 6.62, 3.7))
    add_badge(slide, "现有架构图", Box(5.66, 5.86, 1.08, 0.28), PALETTE["sage"], PALETTE["ink"], font_size=10.5)
    add_text(slide, "如果你想把这页做得更像 keynote，建议后续换成单独的信息图版本。", Box(6.94, 5.82, 4.9, 0.3), 11.8, PALETTE["muted"])
    add_footer(slide, 8, "System Overview")


def build_system_stats(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["warm"])
    add_title_block(slide, "三、体量", "它不是一个 prompt 试玩，而是一个可治理的双层系统")
    stats = [
        ("22,985", "nanobot/ Python 行数", "84 个 Python 文件，上游轻量框架本体。", PALETTE["paper"]),
        ("14,071", "ava/ Python 行数", "80 个 Python 文件，sidecar 扩展层。", PALETTE["paper"]),
        ("13", "Patch 模块", "来自 `.specanchor/modules/module-index.md`。", PALETTE["appendix"]),
        ("8", "自定义工具", "由 `tools_patch` 注入，含 codex / page_agent / gateway_control。", PALETTE["appendix"]),
    ]
    positions = [Box(0.95, 1.9, 2.85, 1.78), Box(4.0, 1.9, 2.85, 1.78), Box(7.05, 1.9, 2.85, 1.78), Box(10.1, 1.9, 2.25, 1.78)]
    for (value, label, note, fill), box in zip(stats, positions, strict=True):
        add_stat_card(slide, value, label, note, box, fill_hex=fill)
    add_card_with_title(
        slide,
        "这组数字想说明什么",
        Box(0.95, 4.18, 5.2, 1.88),
        fill_hex=PALETTE["paper"],
        body_items=[
            "这不是只靠聊天记录维持的“临时 AI 工作流”，而是有清晰目录、运行时、测试和治理面的系统。",
            "真正的价值不是代码多，而是结构分层清楚：上游保持轻，sidecar 承担实验和治理。",
        ],
        body_size=14.2,
    )
    add_card_with_title(
        slide,
        "你可以直接引用的硬证据",
        Box(6.5, 4.18, 5.85, 1.88),
        fill_hex=PALETTE["appendix"],
        body_items=[
            "`.specanchor/global/architecture.md`",
            "`.specanchor/modules/module-index.md`",
            "`.specanchor/patch_map.md`",
            "`ava/agent/bg_tasks.py` / `ava/console/ui_build.py` / `ava/patches/loop_patch.py`",
        ],
        body_size=13.8,
    )
    add_footer(slide, 9, "Repo Facts")


def build_half_loop(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["content"])
    add_title_block(slide, "三、闭环", "带人工监督的研发半闭环")
    flow = [
        ("人类下达任务", "定义目标、边界、是否继续。", Box(0.92, 2.2, 2.0, 1.2), PALETTE["paper"]),
        ("Coding Agent", "Cursor / Claude Code / Codex 产出代码。", Box(3.15, 2.2, 2.15, 1.2), PALETTE["sand"]),
        ("后台任务层", "bg_tasks 持久化 + 通知 + 恢复。", Box(5.58, 2.2, 2.15, 1.2), PALETTE["appendix"]),
        ("自动续跑", "auto_continue 把结果注回原会话。", Box(8.02, 2.2, 2.2, 1.2), PALETTE["paper"]),
        ("自动补救", "auto rebuild / lifecycle status / restart。", Box(10.48, 2.2, 2.0, 1.2), PALETTE["sage"]),
    ]
    for idx, (title, body, box, fill) in enumerate(flow):
        add_flow_box(slide, title, body, box, fill)
        if idx < len(flow) - 1:
            add_arrow(slide, box.x + box.w + 0.1, 2.55, 0.36, 0.42, PALETTE["clay"])
    add_round_rect(slide, Box(1.2, 4.25, 10.95, 1.42), PALETTE["paper"], line_hex=PALETTE["line"])
    add_text(slide, "为什么叫“半闭环”", Box(1.45, 4.42, 2.1, 0.25), 16, PALETTE["ink"], bold=True)
    add_text(
        slide,
        "开发 loop 和 runtime loop 已经能自己续起来，但 release / PR / deploy 仍然留给人类拍板。\n所以人不是被拿掉了，而是从“亲手做所有事”变成“决定何时放行”。",
        Box(1.46, 4.78, 10.2, 0.62),
        14,
        PALETTE["ink"],
    )
    add_footer(slide, 10, "Semi-Closed Loop")


def build_story_slide(prs: Presentation, bgs: dict[str, Path], page_no: int, title: str, symptom: Sequence[str], why: Sequence[str], fix: Sequence[str], mapping: str, fill_hex: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["content"])
    add_title_block(slide, "四、真实问题", title)
    add_card_with_title(slide, "症状", Box(0.9, 1.86, 4.1, 1.9), fill_hex=PALETTE["paper"], body_items=list(symptom), body_size=14.2)
    add_card_with_title(slide, "为什么会这样", Box(0.9, 4.05, 4.1, 2.0), fill_hex=PALETTE["appendix"], body_items=list(why), body_size=13.8)
    add_card_with_title(slide, "怎么修", Box(5.35, 1.86, 6.95, 4.2), fill_hex=fill_hex, body_items=list(fix), body_size=14.1)
    add_badge(slide, mapping, Box(5.56, 5.73, 3.9, 0.33), PALETTE["clay"], PALETTE["paper"], font_size=12)
    add_footer(slide, page_no, "Three Stories")


def build_layers(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["warm"])
    add_title_block(slide, "四、总结图", "约束不是一句提示词，而是分层设计")
    layers = [
        ("Layer 4 · CI / Guardrail Tests", "最终兜底：tests + guardrails + hook 检查", PALETTE["paper"], 1.55),
        ("Layer 3 · Post-task Hooks", "自动补救：auto rebuild / continuation / lifecycle", PALETTE["appendix"], 2.4),
        ("Layer 2 · Tool Description + Spec", "影响 AI 决策：委托策略、模块契约、SpecAnchor", PALETTE["sand"], 3.25),
        ("Layer 1 · AGENTS.md / CLAUDE.md", "最外层说明文档：告诉 AI 默认边界", PALETTE["paper"], 4.1),
    ]
    for title, body, fill, y in layers:
        add_round_rect(slide, Box(1.25, y, 10.9, 0.72), fill, line_hex=PALETTE["line"])
        add_text(slide, title, Box(1.5, y + 0.13, 4.5, 0.22), 16, PALETTE["ink"], bold=True)
        add_text(slide, body, Box(5.55, y + 0.14, 5.9, 0.22), 13.2, PALETTE["ink"])
    add_text(
        slide,
        "每一层都不需要单独 100% 可靠。多层叠加 + sidecar 边界，才能把“AI 会犯错”变成“AI 犯错时成本可控”。",
        Box(1.28, 5.45, 10.7, 0.52),
        15,
        PALETTE["ink"],
        align=PP_ALIGN.CENTER,
        bold=True,
    )
    add_footer(slide, 14, "Layered Guardrails")


def build_personal_edge(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["content"])
    add_title_block(slide, "五、个人优势", "我不是只会踩坑，而是把坑变成了工程能力")
    cards = [
        (
            "架构判断力",
            [
                "选择 sidecar 而不是 fork，把上游同步的长期成本从结构上压低。",
                "知道什么放 patch、什么放 fork、什么该回提 upstream。",
            ],
            PALETTE["paper"],
        ),
        (
            "治理意识",
            [
                "不满足于“写几句提示词”，而是把 prompt 失效处变成测试、hook、运行时控制面。",
                "用 SpecAnchor 把经验沉淀成可加载、可追踪的团队契约。",
            ],
            PALETTE["appendix"],
        ),
        (
            "闭环落地能力",
            [
                "不是只指出问题，而是把问题写成 continuation、auto rebuild、lifecycle control。",
                "把“AI 改 AI”从概念推进到可运行的半闭环。",
            ],
            PALETTE["sand"],
        ),
    ]
    x = 0.9
    for title, body, fill in cards:
        add_card_with_title(slide, title, Box(x, 2.12, 3.96, 3.55), fill_hex=fill, body_items=body, body_size=14.2)
        x += 4.18
    add_footer(slide, 15, "Personal Edge")


def build_patterns(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["warm"])
    add_title_block(slide, "六、可复用模式", "你今天就能带走的三个模式")
    patterns = [
        ("模式 1：Sidecar 扩展", ["依赖开源项目但要深度定制？优先旁路扩展，不要直接 fork。", "前提是上游有足够 patch 点。"], PALETTE["paper"]),
        ("模式 2：软硬护栏组合", ["软护栏决定 AI 想做什么，硬护栏决定 AI 能做什么。", "提示词、Spec、测试、CI、Hook 必须一起上。"], PALETTE["appendix"]),
        ("模式 3：Post-task 自动化", ["build、lint、deploy 这类固定后续动作，不要赌 AI 会记得。", "把它变成系统级 hook。"], PALETTE["sand"]),
    ]
    x = 0.88
    for title, body, fill in patterns:
        add_card_with_title(slide, title, Box(x, 2.08, 3.95, 3.45), fill_hex=fill, body_items=body, body_size=14.2)
        x += 4.18
    add_footer(slide, 16, "Reusable Patterns")


def build_qa(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["dark"])
    add_title_block(slide, "七、Q&A", "预设问题：提前把最容易被追问的点讲清楚", dark=True)
    left = [
        "Monkey Patch 会不会很脆弱？\n会，所以每个 patch 都有 guard 和最低验证。",
        "小项目值得上这套吗？\n如果持续迭代且频繁用 AI，就值得；否则从 AGENTS.md 起步。",
        "能不能让 AI 自己设计约束？\n可以，但人仍要 review 边界情况。",
    ]
    right = [
        "Token 成本会不会更高？\n会涨上下文成本，但通常更少返工，净效果往往更省。",
        "Nanobot 和 OpenClaw 的差别？\n不是谁更强，而是谁更适合做可控实验。",
        "SpecAnchor 和 Spec Kit / OpenSpec / Kiro？\n前者偏治理，后者偏 authoring / workflow。",
    ]
    add_card_with_title(slide, "常见追问 A", Box(0.92, 2.0, 5.75, 4.35), fill_hex=PALETTE["appendix"], body_items=left, body_size=13.8)
    add_card_with_title(slide, "常见追问 B", Box(6.95, 2.0, 5.45, 4.35), fill_hex=PALETTE["paper"], body_items=right, body_size=13.8)
    add_footer(slide, 17, "Q&A", dark=True)


def build_evidence(prs: Presentation, bgs: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["content"])
    add_title_block(slide, "附录", "仓库里已经有的硬证据入口")
    evidence = [
        ("架构规范", "`.specanchor/global/architecture.md`", "用来证明 sidecar / patch / fork 的设计原则。"),
        ("模块覆盖", "`.specanchor/modules/module-index.md`", "13 个 patch 模块与功能模块的总索引。"),
        ("热区治理", "`.specanchor/patch_map.md`", "说明哪些 patch 是 merge hot zone。"),
        ("后台状态", "`ava/agent/bg_tasks.py`", "支撑 continuation、恢复与通知。"),
        ("前端新鲜度", "`ava/console/ui_build.py`", "支撑 auto rebuild 检测。"),
        ("关键提交", "`bd4be432` / `194ccd43` / `1f51f458`", "可支撑故事 2、3 与 lifecycle 的演讲证据。"),
    ]
    x = 0.9
    y = 1.95
    for idx, (title, path_text, note) in enumerate(evidence):
        add_card_with_title(
            slide,
            title,
            Box(x, y, 3.95, 1.45),
            fill_hex=PALETTE["paper"] if idx % 2 == 0 else PALETTE["appendix"],
            body_items=[path_text, note],
            body_size=12.7,
        )
        if idx % 3 == 2:
            x = 0.9
            y += 1.7
        else:
            x += 4.18
    add_footer(slide, 18, "Evidence")


def build_placeholder_slide(prs: Presentation, bgs: dict[str, Path], page_no: int, title: str, tag: str, purpose: str, requirements: Sequence[str], ref_text: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bgs["appendix"])
    add_title_block(slide, "待补素材页", title)
    add_badge(slide, tag, Box(0.9, 1.56, 1.24, 0.32), PALETTE["clay"], PALETTE["paper"], font_size=11.5)
    add_card_with_title(
        slide,
        "这页的用途",
        Box(0.92, 1.95, 3.45, 1.7),
        fill_hex=PALETTE["paper"],
        body_items=[purpose],
        body_size=14.2,
    )
    add_card_with_title(
        slide,
        "请你后续补什么",
        Box(4.65, 1.95, 7.62, 3.75),
        fill_hex=PALETTE["paper"],
        body_items=list(requirements),
        body_size=13.6,
    )
    add_text(slide, ref_text, Box(0.95, 6.3, 11.2, 0.32), 11.8, PALETTE["muted"])
    add_footer(slide, page_no, "Appendix Placeholder")


def build_presentation() -> None:
    backgrounds = ensure_backgrounds()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.title = "怎么让 AI 帮你改代码，但不把项目改坏"
    prs.core_properties.subject = "Nanobot/Ava 技术分享"
    prs.core_properties.keywords = "Nanobot,Ava,Harness Engineering,SpecAnchor,AI Coding"

    build_cover(prs, backgrounds)
    build_core_message(prs, backgrounds)
    build_openclaw(prs, backgrounds)
    build_sources(prs, backgrounds)
    build_harness_five(prs, backgrounds)
    build_ava_mapping(prs, backgrounds)
    build_specanchor_role(prs, backgrounds)
    build_nano_overview(prs, backgrounds)
    build_system_stats(prs, backgrounds)
    build_half_loop(prs, backgrounds)
    build_story_slide(
        prs,
        backgrounds,
        11,
        "故事 1：改错地方",
        symptom=["AI 为了“更省事”，越过边界直接改了 `nanobot/`。"],
        why=[
            "规则只写在提示词里，AI 仍会在“效率”和“规则”之间自己权衡。",
            "它有能力违背规范，所以规范不能只存在于文档里。",
        ],
        fix=[
            "Layer 1：加强 AGENTS / CLAUDE / tool description 的措辞。",
            "Layer 2：guardrail 测试自动检查是否误改 `nanobot/`。",
            "Layer 3：sidecar 架构把默认扩展路径收束到 `ava/`，让越界更容易被拦截。",
        ],
        mapping="Harness 映射：Instructions + Verification + Scope",
        fill_hex=PALETTE["sand"],
    )
    build_story_slide(
        prs,
        backgrounds,
        12,
        "故事 2：改完不收尾",
        symptom=["后台任务完成了，也通知了我，但 agent 没有继续推进下一步。"],
        why=[
            "完成回调只负责通知，没有把结果重新注回原会话。",
            "等于员工汇报了结果，系统却没有安排后续动作。",
        ],
        fix=[
            "补 `auto_continue`：任务完成后自动把结果注回原会话。",
            "加 continuation budget：每个会话最多自动续跑 5 次，避免死循环。",
            "用户发新消息时重置预算，让人仍然是最终方向控制器。",
        ],
        mapping="Harness 映射：Lifecycle + State",
        fill_hex=PALETTE["appendix"],
    )
    build_story_slide(
        prs,
        backgrounds,
        13,
        "故事 3：改完看不到",
        symptom=["前端 TypeScript 已改，AI 也说完成了，但真实页面没有变化。"],
        why=[
            "AI 不一定记得还要 `npm run build`。",
            "这正是 harness 里最危险的错觉：confidence ≠ correctness。",
        ],
        fix=[
            "在 post-task hook 里检测 console-ui 产物是否过期。",
            "源码 mtime 大于 dist mtime 时自动 rebuild。",
            "不要依赖 AI 自己想起“还差一步”。",
        ],
        mapping="Harness 映射：Verification（产物新鲜度检测）",
        fill_hex=PALETTE["paper"],
    )
    build_layers(prs, backgrounds)
    build_personal_edge(prs, backgrounds)
    build_patterns(prs, backgrounds)
    build_qa(prs, backgrounds)
    build_evidence(prs, backgrounds)
    build_placeholder_slide(
        prs,
        backgrounds,
        19,
        "[待补] OpenClaw 视觉 Hook",
        "截图",
        "用于替换第 3 页右侧占位区，强化“很火，但也很危险”的第一印象。",
        [
            "一张 16:9 横向截图，优先包含 repo 星标数或媒体标题。",
            "如果能同屏出现产品界面或相关新闻摘要更好。",
            "避免纯文字长截图，优先“热度”一眼可见。",
        ],
        "更细的截图要求见 `docs/tech-sharing-outline-assets.md`。",
    )
    build_placeholder_slide(
        prs,
        backgrounds,
        20,
        "[待补] Ava 架构全景信息图",
        "信息图",
        "建议替换第 8 页的现有架构图，让系统全貌更适合演讲呈现。",
        [
            "左侧是 Nanobot 核心，右侧是 Ava sidecar，底部是 console / tools / storage / runtime。",
            "请保留“不是 fork，是 sidecar”这个信息焦点。",
            "推荐用明亮柔和、有人文温度的插图式信息图。",
        ],
        "Nano Banana 提示词已写入 `docs/tech-sharing-outline-assets.md`。",
    )
    build_placeholder_slide(
        prs,
        backgrounds,
        21,
        "[待补] 1-2 分钟 Demo 录屏",
        "视频",
        "这页用于在演讲时播放或截图拆帧，展示“带人工监督的 AI 改 AI”真实链路。",
        [
            "建议一镜到底：任务下发 → 后台执行 → 通知 → auto continue → 页面更新。",
            "画面至少同时出现一个沟通界面和一个开发/控制界面。",
            "时长控制在 60-120 秒，字幕可后补。",
        ],
        "建议镜头脚本见 `docs/tech-sharing-outline-assets.md`。",
    )
    build_placeholder_slide(
        prs,
        backgrounds,
        22,
        "[待补] 三个故事的证据截图",
        "截图",
        "用于补强第 11-13 页，让三个故事从“说得通”变成“看得见”。",
        [
            "故事 1：误改 `nanobot/` 的 diff 或 guardrail 拦截截图。",
            "故事 2：任务完成通知 + auto continue 继续执行的连续截图。",
            "故事 3：前端未生效 / 重建后生效的对比截图。",
        ],
        "建议截图组合方式见 `docs/tech-sharing-outline-assets.md`。",
    )
    build_placeholder_slide(
        prs,
        backgrounds,
        23,
        "[待补] Commit / Diff 硬证据",
        "代码截图",
        "用于第 15 页之前或附录中插入真实研发证据，增强说服力。",
        [
            "优先准备 `bd4be432`、`194ccd43`、`1f51f458` 这类能对应故事或闭环的提交。",
            "再补一组 `loop_patch.py` 的关键 diff，说明不是纸上谈兵。",
            "若要讲“AI 误改上游”，请从本地历史中挑一个代表性案例截图。",
        ],
        "如何挑 commit、如何截 diff，见 `docs/tech-sharing-outline-assets.md`。",
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
